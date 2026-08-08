"""Making documents — text, markdown, spreadsheets, slides.

The wish-list item is "create, edit and manage documents, presentations,
spreadsheets and media". Most of that is one honest observation: the formats
worth writing by voice are the ones that are just text. Markdown is a
document, CSV is a spreadsheet, and both survive being edited by hand
afterwards, which is what actually happens.

So the native formats here are text, markdown and CSV, and they have no
dependencies at all. The binary formats — .docx, .xlsx, .pptx — are real
requests too, and they go through optional libraries that degrade to a clear
message when absent, the same way `reading.py` handles PDFs. On a 2 GB Pi
those libraries may deliberately never be installed; the answer "I can write
you a markdown file but not a Word document on this device" is a useful one.

Two rules everything here obeys:

  * **Never overwrite without being told to.** A voice assistant asked to
    "write up the meeting notes" that silently replaces yesterday's notes has
    destroyed something no one can get back. Existing files are refused
    unless `overwrite=True`.
  * **Never write outside the folder it was given.** A filename is treated as
    a filename, not a path — "../../.ssh/authorized_keys" is a name with
    slashes in it, and it gets flattened rather than followed.
"""

from __future__ import annotations

import csv
import io
import logging
import re
from collections.abc import Sequence
from dataclasses import dataclass
from pathlib import Path

log = logging.getLogger("flint.documents")

#: Written documents are small. A voice assistant producing a 50 MB file has
#: misunderstood something.
MAX_CHARS = 2_000_000

_UNSAFE = re.compile(r"[^A-Za-z0-9._ \-()]+")

FORMATS = {
    "txt": "text", "md": "markdown", "markdown": "markdown",
    "csv": "spreadsheet", "tsv": "spreadsheet",
    "docx": "word document", "xlsx": "spreadsheet", "pptx": "presentation",
}


class DocumentError(Exception):
    pass


def safe_name(name: str, default_extension: str = "txt") -> str:
    """A filename, guaranteed to stay in the folder it is joined to.

    Path separators are not honoured — they are stripped. Anything asking for
    a directory is asking for something this module does not do.
    """
    raw = (name or "").strip().replace("\\", "/")
    base = raw.rsplit("/", 1)[-1]          # drop any directory part entirely
    base = _UNSAFE.sub("_", base).strip(". ")
    if not base:
        base = "document"
    if "." not in base:
        base = f"{base}.{default_extension}"
    return base[:120]


def extension_of(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


@dataclass(frozen=True)
class WriteResult:
    ok: bool
    path: str = ""
    detail: str = ""

    def spoken(self) -> str:
        if not self.ok:
            return self.detail
        return f"Saved as {Path(self.path).name}."


def _resolve(folder: str | Path, name: str, extension: str,
             overwrite: bool) -> Path:
    directory = Path(folder).expanduser()
    if not directory.is_dir():
        raise DocumentError(f"there's no folder at {directory}")
    target = directory / safe_name(name, extension)
    if target.exists() and not overwrite:
        raise DocumentError(
            f"{target.name} already exists — say overwrite it if you mean to "
            f"replace it, or give me a different name")
    return target


def _write_text(target: Path, text: str) -> WriteResult:
    if len(text) > MAX_CHARS:
        return WriteResult(False, detail="that's far too long to write out")
    try:
        target.write_text(text, encoding="utf-8")
    except OSError as exc:
        return WriteResult(False, detail=f"couldn't write it: {exc}")
    return WriteResult(True, path=str(target))


def write_document(folder: str | Path, name: str, content: str, *,
                   title: str = "", overwrite: bool = False) -> WriteResult:
    """A text or markdown document. `title` becomes an H1 for markdown."""
    try:
        target = _resolve(folder, name, "md", overwrite)
    except DocumentError as exc:
        return WriteResult(False, detail=str(exc))

    extension = extension_of(target.name)
    if extension == "docx":
        return _write_docx(target, content, title)
    body = content or ""
    if title and extension in ("md", "markdown"):
        body = f"# {title}\n\n{body}"
    return _write_text(target, body)


def write_spreadsheet(folder: str | Path, name: str,
                      rows: Sequence[Sequence], *, headers: Sequence[str] = (),
                      overwrite: bool = False) -> WriteResult:
    """A spreadsheet. CSV natively; .xlsx when openpyxl is available."""
    try:
        target = _resolve(folder, name, "csv", overwrite)
    except DocumentError as exc:
        return WriteResult(False, detail=str(exc))

    table = ([list(headers)] if headers else []) + [list(r) for r in rows]
    if not table:
        return WriteResult(False, detail="there's nothing to put in it")

    if extension_of(target.name) == "xlsx":
        return _write_xlsx(target, table)

    buffer = io.StringIO()
    # lineterminator="\n" on purpose: csv.writer defaults to "\r\n", and the
    # text write then translates the "\n" again, giving "\r\r\n" — a file that
    # opens with a blank line between every row.
    writer = csv.writer(buffer, lineterminator="\n",
                        delimiter="\t" if target.name.endswith(".tsv") else ",")
    writer.writerows(table)
    return _write_text(target, buffer.getvalue())


def write_presentation(folder: str | Path, name: str,
                       slides: Sequence[dict], *, title: str = "",
                       overwrite: bool = False) -> WriteResult:
    """Slides. Markdown (`---` separated) natively; .pptx via python-pptx.

    Markdown slides are the default on purpose: they open anywhere, diff
    properly, and every presentation tool imports them.
    """
    try:
        target = _resolve(folder, name, "md", overwrite)
    except DocumentError as exc:
        return WriteResult(False, detail=str(exc))
    if not slides:
        return WriteResult(False, detail="there are no slides to write")

    if extension_of(target.name) == "pptx":
        return _write_pptx(target, slides, title)

    blocks = []
    if title:
        blocks.append(f"# {title}")
    for slide in slides:
        lines = [f"## {slide.get('title', '')}".rstrip()]
        for bullet in slide.get("bullets", ()) or ():
            lines.append(f"- {bullet}")
        if slide.get("notes"):
            lines.append(f"\n> {slide['notes']}")
        blocks.append("\n".join(lines).strip())
    return _write_text(target, "\n\n---\n\n".join(blocks) + "\n")


# ── binary formats, all optional ─────────────────────────────────────────────
def _write_docx(target: Path, content: str, title: str) -> WriteResult:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return WriteResult(False, detail=(
            "I can't write Word files on this device — python-docx isn't "
            "installed. I can do it as markdown if that's any use."))
    document = DocxDocument()
    if title:
        document.add_heading(title, level=1)
    for paragraph in (content or "").split("\n\n"):
        document.add_paragraph(paragraph)
    try:
        document.save(str(target))
    except OSError as exc:
        return WriteResult(False, detail=f"couldn't write it: {exc}")
    return WriteResult(True, path=str(target))


def _write_xlsx(target: Path, table: Sequence[Sequence]) -> WriteResult:
    try:
        from openpyxl import Workbook
    except ImportError:
        return WriteResult(False, detail=(
            "I can't write Excel files on this device — openpyxl isn't "
            "installed. I can save it as a CSV instead, which opens in Excel."))
    book = Workbook()
    sheet = book.active
    for row in table:
        sheet.append(list(row))
    try:
        book.save(str(target))
    except OSError as exc:
        return WriteResult(False, detail=f"couldn't write it: {exc}")
    return WriteResult(True, path=str(target))


def _write_pptx(target: Path, slides: Sequence[dict], title: str) -> WriteResult:
    try:
        from pptx import Presentation
    except ImportError:
        return WriteResult(False, detail=(
            "I can't write PowerPoint on this device — python-pptx isn't "
            "installed. I can write the slides as markdown instead."))
    deck = Presentation()
    for slide in slides:
        layout = deck.slide_layouts[1]
        added = deck.slides.add_slide(layout)
        added.shapes.title.text = str(slide.get("title", title or ""))
        body = added.placeholders[1].text_frame
        bullets = list(slide.get("bullets", ()) or ())
        if bullets:
            body.text = str(bullets[0])
            for bullet in bullets[1:]:
                body.add_paragraph().text = str(bullet)
    try:
        deck.save(str(target))
    except OSError as exc:
        return WriteResult(False, detail=f"couldn't write it: {exc}")
    return WriteResult(True, path=str(target))


# ── reading back ─────────────────────────────────────────────────────────────
def read_document(path: str | Path, max_chars: int = 20_000) -> str:
    """The text of a document she wrote, so it can be edited or read out."""
    target = Path(path).expanduser()
    if not target.is_file():
        return f"there's no file at {target}"
    if extension_of(target.name) in ("docx", "xlsx", "pptx"):
        return (f"{target.name} is a {FORMATS.get(extension_of(target.name), 'binary')} "
                f"— I can write those but not read them back on this device")
    try:
        return target.read_text(encoding="utf-8", errors="replace")[:max_chars]
    except OSError as exc:
        return f"couldn't read it: {exc}"


def list_documents(folder: str | Path, limit: int = 30) -> list[str]:
    directory = Path(folder).expanduser()
    try:
        found = [p for p in directory.iterdir()
                 if p.is_file() and extension_of(p.name) in FORMATS]
    except OSError:
        return []
    found.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    return [p.name for p in found[:limit]]
